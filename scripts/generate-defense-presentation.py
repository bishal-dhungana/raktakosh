from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Raktakosh_Defense_Proposal_Presentation_Updated.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(198, 47, 52)
PRESENTATION_RED = RGBColor(225, 60, 64)
INK = RGBColor(42, 36, 38)
TEXT = RGBColor(68, 63, 65)
MUTED = RGBColor(113, 106, 108)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(253, 244, 244)

FONT = "Aptos"
DISPLAY = "Aptos Display"


def background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def text(slide, value, x, y, w, h, size=18, color=TEXT, bold=False,
         font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = value
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.font.language_id = 1033
    return box


def bullet_list(slide, values, x, y, w, h, size=18, numbered=False, gap=10):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for i, value in enumerate(values):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        prefix = f"{i + 1}. " if numbered else "• "
        paragraph.text = prefix + value
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(gap)
        paragraph.font.language_id = 1033
    return box


def title(slide, heading, subheading=None):
    text(slide, heading, 0.78, 0.58, 11.75, 0.55, size=32, color=INK, bold=True, font=DISPLAY)
    if subheading:
        text(slide, subheading, 0.80, 1.24, 11.65, 0.34, size=18, color=MUTED)


def section(slide, heading, x, y, w):
    text(slide, heading, x, y, w, 0.34, size=21, color=RED, bold=True, font=DISPLAY)


def fill_shape(slide, shape_type, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    return shape


def diagram_box(slide, heading, body, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.25)
    text(slide, heading, x + 0.12, y + 0.16, w - 0.24, 0.31, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(slide, body, x + 0.14, y + 0.58, w - 0.28, h - 0.70, size=16, color=TEXT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def diagram_arrow(slide, x, y, w, h, direction="right"):
    arrow_type = SHAPE.RIGHT_ARROW if direction == "right" else SHAPE.DOWN_ARROW
    shape = slide.shapes.add_shape(arrow_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.color.rgb = RED


def arrow_label(slide, value, x, y, w):
    text(slide, value, x, y, w, 0.22, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)


# 1. Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = PALE
fill_shape(slide, SHAPE.RECTANGLE, 0, 0, 4.48, 7.5, PRESENTATION_RED)
fill_shape(slide, SHAPE.ROUNDED_RECTANGLE, 0.72, 0.68, 0.72, 0.72, WHITE)
fill_shape(slide, SHAPE.TEAR, 0.91, 0.78, 0.34, 0.48, PRESENTATION_RED)
text(slide, "RAKTAKOSH", 0.72, 1.59, 3.0, 0.34, size=20, color=WHITE, bold=True, font=DISPLAY)
text(slide, "BLOOD COORDINATION PLATFORM", 0.73, 1.97, 3.25, 0.26, size=16, color=WHITE, bold=True)
text(slide, "MINOR PROJECT\nDEFENSE PROPOSAL", 0.72, 3.10, 3.14, 1.14, size=26, color=WHITE, bold=True, font=DISPLAY)
text(slide, "PRJ 360 · Department of Computer Engineering", 0.72, 4.68, 3.20, 0.33, size=16, color=WHITE)
text(slide, "Madan Bhandari College of Engineering\nUrlabari–3, Morang", 0.72, 6.17, 3.10, 0.48, size=16, color=WHITE)
text(slide, "REAL-TIME BLOOD INVENTORY & DONOR MATCHING", 5.18, 0.96, 7.00, 0.30, size=16, color=PRESENTATION_RED, bold=True)
text(slide, "Raktakosh", 5.14, 1.57, 6.6, 0.68, size=38, color=INK, bold=True, font=DISPLAY)
text(slide, "A Real-Time Web System for Blood Inventory\nand Donor Matching", 5.17, 2.33, 7.0, 0.8, size=24, color=TEXT, font=DISPLAY)
fill_shape(slide, SHAPE.RECTANGLE, 5.18, 3.43, 6.75, 0.018, RGBColor(226, 217, 218))
text(slide, "Submitted by", 5.18, 3.78, 2.2, 0.28, size=17, color=MUTED, bold=True)
text(slide, "Sambridhee Dhakal [23070563]\nDeepika Adhikari [23070542]\nNischal Dahal [22070029]", 5.18, 4.14, 5.8, 1.05, size=19, color=INK, bold=True)
text(slide, "Defense presentation", 5.18, 6.46, 4.8, 0.28, size=17, color=PRESENTATION_RED, bold=True)
text(slide, "2026", 11.75, 6.46, 0.8, 0.28, size=17, color=MUTED, align=PP_ALIGN.RIGHT)


# 2. Outline
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Presentation outline")
outline = [
    "Introduction & background",
    "Problem statement",
    "Objectives",
    "Literature review & research gap",
    "Proposed solution & scope",
    "Methodology & system design",
    "Implementation & deployment",
    "Gantt chart & future work",
    "Conclusion",
]
bullet_list(slide, outline, 1.00, 1.82, 8.6, 4.90, size=21, numbered=True, gap=13)


# 3. Introduction
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Introduction & background", "Raktakosh supports responsible blood-service coordination through a simple web platform.")
bullet_list(slide, [
    "Connects people seeking blood, voluntary donors, verified Blood Banks, and administrators.",
    "Helps families locate verified Blood Bank information without repeated manual searching.",
    "Supports private blood requests and facility-led coordination.",
    "Uses consent-first donor communication and privacy-aware workflows.",
    "Clinical decisions remain with licensed Blood Banks.",
], 0.92, 2.00, 11.30, 3.90, size=20, gap=15)


# 4. Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Problem statement", "Blood emergencies can become more difficult when availability and communication are fragmented.")
bullet_list(slide, [
    "Families may call multiple facilities or make public social-media appeals to find blood.",
    "Blood availability information can be difficult to find or confirm before travelling.",
    "Personal details can be exposed when requests are handled through public posts.",
    "Manual coordination may not provide a clear status history or accountable follow-up.",
    "A verified, private, and facility-led coordination flow is needed.",
], 0.92, 2.00, 11.30, 3.90, size=20, gap=15)


# 5. Objectives
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Project objectives")
bullet_list(slide, [
    "Centralize coordination between seekers, donors, verified Blood Banks, and administrators.",
    "Improve visibility of facility-reported blood availability by district and component.",
    "Support private requests, document review, status updates, and facility follow-up.",
    "Maintain consent-based donor coordination when facility inventory is unavailable.",
    "Strengthen accountability through role-based access and audit records.",
], 0.92, 1.88, 11.30, 4.10, size=20, numbered=True, gap=15)


# 6. Literature review
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Literature review & research gap")
section(slide, "Hamro Life Bank", 0.90, 1.84, 5.3)
bullet_list(slide, [
    "Supports emergency blood search through a hotline and online communication channels.",
    "Shows the importance of organized blood-request coordination.",
    "Research gap: the process remains heavily dependent on centralized operators.",
], 0.92, 2.35, 5.55, 2.38, size=18, gap=11)
section(slide, "Nepal Red Cross Society", 6.82, 1.84, 5.3)
bullet_list(slide, [
    "Supports blood-service operations, stock monitoring, and donor record management.",
    "Shows the value of digital inventory and donor information systems.",
    "Research gap: limited end-to-end private requester-to-facility coordination.",
], 6.84, 2.35, 5.50, 2.38, size=18, gap=11)
section(slide, "Raktakosh response", 0.90, 5.18, 5.3)
bullet_list(slide, [
    "A single platform for discovery, requests, facility workflows, donor outreach, and auditability.",
], 0.92, 5.70, 11.15, 0.72, size=19, gap=10)


# 7. Scope
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Proposed solution & scope")
section(slide, "What the platform provides", 0.90, 1.86, 5.3)
bullet_list(slide, [
    "Public Blood Bank discovery by district, group, Rh factor, and component.",
    "Private blood requests with supporting documents for facility review.",
    "Facility inventory updates, case status tracking, and donor outreach workflows.",
    "Responsive English and Nepali public interface support.",
], 0.92, 2.36, 5.52, 3.00, size=18, gap=11)
section(slide, "Important boundaries", 6.82, 1.86, 5.3)
bullet_list(slide, [
    "Does not determine clinical compatibility or medical donor eligibility.",
    "Does not guarantee reservations, laboratory testing, or transfusion availability.",
    "Does not replace the responsibility of participating Blood Banks.",
], 6.84, 2.36, 5.46, 2.56, size=18, gap=11)
text(slide, "Every clinical decision remains with the licensed Blood Bank.", 6.84, 5.45, 5.25, 0.50, size=20, color=RED, bold=True)


# 8. Methodology and design
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Methodology & system design")
section(slide, "Iterative Agile methodology", 0.90, 1.84, 5.3)
bullet_list(slide, [
    "Requirement analysis: identify users, problems, and project scope.",
    "System design: plan the data structure, modules, and user workflows.",
    "Development: build frontend, backend, database, and core workflows.",
    "Testing and review: refine features using feedback and verification.",
], 0.92, 2.35, 5.50, 3.20, size=18, numbered=True, gap=12)
section(slide, "Three-tier system design", 6.82, 1.84, 5.3)
bullet_list(slide, [
    "Presentation layer: React, TypeScript, and Vite for responsive user interfaces.",
    "Application layer: Node.js and Express for validation, workflows, and access control.",
    "Data layer: TiDB Cloud or MySQL for users, inventory, requests, and audit data.",
    "Deployment: Vercel frontend and Render API hosting.",
], 6.84, 2.35, 5.46, 3.20, size=18, numbered=True, gap=12)


# 9. Roles and modules
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Users, roles & key modules")
section(slide, "Users and responsibilities", 0.90, 1.84, 5.3)
bullet_list(slide, [
    "Guest or requester: search public information and submit a private request.",
    "Donor: manage consent and availability, then respond to controlled outreach.",
    "Facility team: update availability, review cases, and coordinate next steps.",
    "Platform administrator: verify facilities and review governance activity.",
], 0.92, 2.35, 5.50, 3.15, size=18, gap=11)
section(slide, "Key modules", 6.82, 1.84, 5.3)
bullet_list(slide, [
    "Public discovery",
    "Request coordination",
    "Donor management",
    "Facility inventory and request review",
    "Governance and audit",
], 6.84, 2.35, 5.48, 3.15, size=18, gap=11)


# 10. Workflow
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "End-to-end coordination workflow")
bullet_list(slide, [
    "A user searches for a Blood Bank using district and blood-component information.",
    "The requester submits a private request with the required supporting document.",
    "Authorized facility staff review the request and reported availability.",
    "If inventory is available, the facility coordinates the next safe operational step.",
    "If inventory is unavailable, the facility may start controlled consent-based donor outreach.",
    "Status updates and audit records are maintained throughout the workflow.",
], 0.92, 1.92, 11.32, 4.68, size=20, numbered=True, gap=15)


# 11. Technology
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Technology stack & deployment")
bullet_list(slide, [
    "Frontend: React, TypeScript, and Vite for a responsive public website and role-aware workspaces.",
    "Backend: Node.js and Express for validation, workflow rules, sessions, and API endpoints.",
    "Database: TiDB Cloud or MySQL for users, inventory, requests, policies, and audit records.",
    "Deployment: Vercel hosts the frontend and Render hosts the API.",
    "Security: server-side role checks, CSRF protection, rate limits, and controlled document access.",
], 0.92, 1.92, 11.32, 4.20, size=20, gap=15)


# 12. Current website
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Current website implementation")
bullet_list(slide, [
    "The deployed public website is available for Blood Bank discovery and public information.",
    "Private request workflow and facility inventory workflow have been implemented.",
    "Facility review and controlled donor coordination features are included in the current system.",
    "User authentication and completion of the user-side account flow are still in progress.",
    "Guest search remains available without creating an account.",
], 0.92, 1.92, 11.32, 4.12, size=20, gap=15)
text(slide, "Current status: public discovery, request pathway, and facility workflow are available; user authentication is in progress.", 0.92, 6.20, 11.18, 0.52, size=18, color=RED, bold=True)


# 13. System data flow diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "System data flow diagram", "This diagram shows how information moves through the Raktakosh platform.")

diagram_box(slide, "Platform users", "Guest / requester\nDonor\nBlood Bank staff", 0.73, 2.45, 2.30, 1.58)
diagram_box(slide, "Web app", "Public search\nPrivate requests\nRole-based dashboards", 3.62, 2.45, 2.43, 1.58)
diagram_box(slide, "API & workflow", "Validation\nAccess control\nRequest and outreach workflow", 6.64, 2.45, 2.56, 1.58)
diagram_box(slide, "Database", "Users and facilities\nInventory and requests\nDocuments and audit records", 9.78, 2.45, 2.63, 1.58)
diagram_box(slide, "System outputs", "Availability results · request status updates · facility coordination · controlled donor outreach", 3.66, 5.13, 5.66, 1.23)

diagram_arrow(slide, 3.08, 3.07, 0.42, 0.22)
diagram_arrow(slide, 6.10, 3.07, 0.42, 0.22)
diagram_arrow(slide, 9.24, 3.07, 0.42, 0.22)
diagram_arrow(slide, 7.75, 4.13, 0.28, 0.67, direction="down")

text(slide, "Staff update inventory and review requests through the same web app. The API reads and writes the data, then returns results and status updates to users.", 0.84, 6.72, 11.58, 0.42, size=16, color=MUTED, align=PP_ALIGN.CENTER)


# 14. Gantt chart
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Gantt chart: project plan")
section(slide, "Completed and current work", 0.90, 1.84, 5.3)
bullet_list(slide, [
    "May: problem study and literature review.",
    "June: requirements analysis and system design.",
    "June–July: UI and public website development.",
    "July: backend, database, and workflow implementation.",
    "July: user authentication completion is in progress.",
    "August: testing, refinement, and defense preparation.",
], 0.92, 2.35, 6.35, 4.15, size=19, numbered=True, gap=12)
section(slide, "Timeline summary", 8.00, 1.84, 4.0)
bullet_list(slide, [
    "Completed: research, design, UI, backend, and core workflows.",
    "In progress: user authentication.",
    "Planned: final testing, refinement, and presentation defense.",
], 8.02, 2.35, 4.08, 2.82, size=18, gap=13)


# 15. Future work
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Future work")
bullet_list(slide, [
    "Complete user authentication and account verification.",
    "Finish end-to-end testing for request, inventory, and role-based workflows.",
    "Add SMS or email notification integration for important updates.",
    "Conduct usability testing with donors, families, and facility staff.",
    "Onboard additional verified Blood Banks and improve coverage.",
    "Improve Nepali content, accessibility, monitoring, and reporting.",
], 0.92, 1.92, 11.32, 4.72, size=20, numbered=True, gap=15)


# 16. Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, "Conclusion")
text(slide, "Raktakosh brings structure to a time-sensitive coordination problem.", 0.92, 1.78, 11.12, 0.74, size=31, color=INK, bold=True, font=DISPLAY)
bullet_list(slide, [
    "Provides one platform for verified discovery, private requests, and facility-led coordination.",
    "Improves transparency through reported availability, workflow status, and audit records.",
    "Uses consent-first donor coordination while keeping clinical decisions with Blood Banks.",
    "The current prototype is deployed, with user authentication identified as the remaining in-progress work.",
], 0.92, 3.00, 11.28, 2.88, size=20, gap=15)
text(slide, "Thank you", 0.92, 6.35, 4.00, 0.42, size=24, color=RED, bold=True, font=DISPLAY)
text(slide, "Questions and discussion", 0.92, 6.79, 5.20, 0.30, size=18, color=MUTED)


prs.core_properties.title = "Raktakosh — Minor Project Defense Proposal"
prs.core_properties.subject = "Raktakosh blood coordination platform"
prs.core_properties.author = "Sambridhee Dhakal, Deepika Adhikari, Nischal Dahal"
prs.save(OUT)
print(f"Created {OUT}")
